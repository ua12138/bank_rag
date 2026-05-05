from __future__ import annotations

import pathlib
import re
import zipfile

from hz_bank_rag.ingestion.multimodal import image_bytes_to_text, image_to_text


TEXT_SUFFIXES = {".txt", ".md", ".log", ".csv"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
PDF_SUFFIXES = {".pdf"}
WORD_SUFFIXES = {".docx", ".doc"}
PPT_SUFFIXES = {".pptx", ".ppt"}


def _strip_xml_tags(text: str) -> str:
    return re.sub(r"<[^>]+>", " ", text).strip()


def parse_document(file_path: str) -> str:
    path = pathlib.Path(file_path)
    suffix = path.suffix.lower()

    if suffix in TEXT_SUFFIXES:
        return _parse_text(path)
    if suffix in PDF_SUFFIXES:
        return _parse_pdf(path)
    if suffix in WORD_SUFFIXES:
        return _parse_word(path)
    if suffix in PPT_SUFFIXES:
        return _parse_ppt(path)
    if suffix in IMAGE_SUFFIXES:
        return image_to_text(str(path))
    return _parse_text(path)


def _parse_text(path: pathlib.Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _parse_pdf(path: pathlib.Path) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        parts: list[str] = []

        for page_index, page in enumerate(reader.pages, start=1):
            txt = (page.extract_text() or "").strip()
            if txt:
                parts.append(f"[pdf-page-{page_index}-text]\n{txt}")

            page_images = getattr(page, "images", []) or []
            image_texts: list[str] = []
            for image_index, img in enumerate(page_images, start=1):
                try:
                    data = getattr(img, "data", b"") or b""
                    name = getattr(img, "name", "") or ""
                    suffix = pathlib.Path(name).suffix or ".png"
                    if data:
                        ocr = image_bytes_to_text(data, suffix=suffix)
                        if ocr:
                            image_texts.append(f"[pdf-page-{page_index}-image-{image_index}-ocr]\n{ocr}")
                except Exception:
                    continue

            if image_texts:
                parts.append("\n".join(image_texts))

        text = "\n\n".join([part for part in parts if part.strip()])
        if text:
            return text
    except Exception:
        pass

    return f"[pdf-parse-fallback] file={path.name}, no text extracted."


def _parse_word(path: pathlib.Path) -> str:
    if path.suffix.lower() == ".docx":
        return _parse_docx(path)
    return _parse_doc(path)


def _parse_docx(path: pathlib.Path) -> str:
    try:
        import docx

        document = docx.Document(str(path))
        lines: list[str] = []

        for p in document.paragraphs:
            text = _strip_xml_tags(p.text or "")
            if text:
                lines.append(text)

        for table in document.tables:
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = _strip_xml_tags(cell.text or "")
                    if cell_text:
                        cells.append(cell_text.replace("\n", " "))
                if cells:
                    lines.append(" | ".join(cells))

        text = "\n".join(lines).strip()
        if text:
            return text
    except Exception:
        pass

    fallback = _parse_docx_xml_fallback(path)
    if fallback:
        return fallback

    return f"[docx-parse-fallback] file={path.name}, no text extracted."


def _parse_docx_xml_fallback(path: pathlib.Path) -> str:
    try:
        with zipfile.ZipFile(path, "r") as zf:
            if "word/document.xml" not in zf.namelist():
                return ""
            xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")

        texts = re.findall(r"<w:t[^>]*>(.*?)</w:t>", xml)
        cleaned = [_strip_xml_tags(t) for t in texts if t and _strip_xml_tags(t)]
        return "\n".join(cleaned).strip()
    except Exception:
        return ""


def _parse_doc(path: pathlib.Path) -> str:
    try:
        import win32com.client  # type: ignore

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(path.resolve()))
        text = doc.Content.Text
        doc.Close(False)
        word.Quit()
        text = (text or "").strip()
        if text:
            return text
    except Exception:
        pass

    return f"[doc-parse-fallback] file={path.name}, suggest converting to .docx first."


def _parse_ppt(path: pathlib.Path) -> str:
    if path.suffix.lower() == ".pptx":
        return _parse_pptx(path)
    return _parse_ppt_legacy(path)


def _parse_pptx(path: pathlib.Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = (getattr(shape, "text", "") or "").strip()
                if text:
                    lines.append(text)
        content = "\n".join(lines).strip()
        if content:
            return content
    except Exception:
        pass

    return f"[pptx-parse-fallback] file={path.name}, no text extracted."


def _parse_ppt_legacy(path: pathlib.Path) -> str:
    try:
        import win32com.client  # type: ignore

        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1
        presentation = app.Presentations.Open(str(path.resolve()), WithWindow=False)

        lines: list[str] = []
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if getattr(shape, "HasTextFrame", False) and shape.TextFrame.HasText:
                    text = shape.TextFrame.TextRange.Text.strip()
                    if text:
                        lines.append(text)

        presentation.Close()
        app.Quit()

        content = "\n".join(lines).strip()
        if content:
            return content
    except Exception:
        pass

    return f"[ppt-parse-fallback] file={path.name}, suggest converting to .pptx first."
